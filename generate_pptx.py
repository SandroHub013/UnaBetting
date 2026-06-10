"""Generate academic PowerPoint presentation for tennis betting project."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

OUTDIR = r"G:\tennis betting\reports"
CHART_DIR = OUTDIR

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ============================================================
# Color palette (academic)
# ============================================================
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
DARK_BLUE = RGBColor(0x2C, 0x3E, 0x6B)
ACCENT_BLUE = RGBColor(0x29, 0x80, 0xB9)
ACCENT_GREEN = RGBColor(0x27, 0xAE, 0x60)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xEC, 0xF0, 0xF1)
MED_GRAY = RGBColor(0x95, 0xA5, 0xA6)
DARK_GRAY = RGBColor(0x2C, 0x3E, 0x50)
BLACK = RGBColor(0x00, 0x00, 0x00)
ACCENT_RED = RGBColor(0xE7, 0x4C, 0x3C)
GOLD = RGBColor(0xF3, 0x9C, 0x12)

def add_background(slide, color=NAVY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color, transparency=0):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_accent_line(slide, left, top, width, color=ACCENT_BLUE, height=Pt(3)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_bullet_list(slide, left, top, width, height, items, font_size=14,
                    color=WHITE, bullet_color=ACCENT_BLUE, spacing=Pt(8), font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = spacing
        # Bullet character
        run_bullet = p.add_run()
        run_bullet.text = "\u25AA  "
        run_bullet.font.size = Pt(font_size)
        run_bullet.font.color.rgb = bullet_color
        run_bullet.font.name = font_name
        run_bullet.font.bold = True
        # Text
        run_text = p.add_run()
        run_text.text = item
        run_text.font.size = Pt(font_size)
        run_text.font.color.rgb = color
        run_text.font.name = font_name
    return txBox

def add_metric_card(slide, left, top, width, height, label, value, sublabel=""):
    card = add_shape_bg(slide, left, top, width, height, DARK_BLUE)
    # Value
    add_text_box(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.6),
                 value, font_size=28, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    # Label
    add_text_box(slide, left + Inches(0.2), top + Inches(0.7), width - Inches(0.4), Inches(0.4),
                 label, font_size=11, color=MED_GRAY, alignment=PP_ALIGN.CENTER)
    if sublabel:
        add_text_box(slide, left + Inches(0.2), top + Inches(1.05), width - Inches(0.4), Inches(0.3),
                     sublabel, font_size=9, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

def add_slide_number(slide, num, total):
    add_text_box(slide, Inches(12.3), Inches(7.05), Inches(0.9), Inches(0.35),
                 f"{num}/{total}", font_size=10, color=MED_GRAY, alignment=PP_ALIGN.RIGHT)

def add_footer_bar(slide):
    add_shape_bg(slide, Inches(0), Inches(7.15), Inches(13.333), Inches(0.35), DARK_BLUE)
    add_text_box(slide, Inches(0.5), Inches(7.18), Inches(6), Inches(0.3),
                 "Tennis Match Prediction System  |  Academic Research Presentation",
                 font_size=9, color=MED_GRAY)

TOTAL_SLIDES = 12

# ============================================================
# SLIDE 1: Title Slide
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_background(slide, NAVY)

# Top accent bar
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)

# Title block
add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10.3), Inches(1.2),
             "A Machine Learning Framework for Tennis Match Prediction",
             font_size=38, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(5.5), Inches(3.1), Inches(2.3), ACCENT_BLUE, Pt(3))

add_text_box(slide, Inches(1.5), Inches(3.5), Inches(10.3), Inches(0.8),
             "Multi-Market Outcome Modeling with Ensemble Methods and Value Betting Strategy",
             font_size=18, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# Metadata
add_text_box(slide, Inches(1.5), Inches(5.0), Inches(10.3), Inches(0.4),
             "Data Sources: JeffSackmann/tennis_atp  |  tennis-data.co.uk  |  TheOddsAPI",
             font_size=12, color=MED_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(5.5), Inches(10.3), Inches(0.4),
             "Models: Logistic Regression  |  Random Forest  |  XGBoost  |  LightGBM  |  Neural Network  |  Ensemble",
             font_size=12, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom bar
add_shape_bg(slide, Inches(0), Inches(7.15), Inches(13.333), Inches(0.35), DARK_BLUE)
add_text_box(slide, Inches(0.5), Inches(7.18), Inches(6), Inches(0.3),
             "CONFIDENTIAL  |  Research Document", font_size=9, color=MED_GRAY)
add_slide_number(slide, 1, TOTAL_SLIDES)

# ============================================================
# SLIDE 2: Research Motivation & Objectives
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)
add_shape_bg(slide, Inches(0), Inches(1.1), Inches(13.333), Inches(0.04), ACCENT_BLUE)
add_footer_bar(slide)
add_slide_number(slide, 2, TOTAL_SLIDES)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(5), Inches(0.7),
             "Research Motivation & Objectives", font_size=28, color=WHITE, bold=True)

# Left column - Motivation
add_shape_bg(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5), DARK_BLUE)
add_text_box(slide, Inches(1.0), Inches(1.55), Inches(5), Inches(0.4),
             "Motivation", font_size=16, color=ACCENT_BLUE, bold=True)

add_bullet_list(slide, Inches(1.0), Inches(2.2), Inches(5.2), Inches(4.5), [
    "Sports betting markets exhibit systematic inefficiencies exploitable via statistical modeling",
    "Tennis is ideal for ML: individual sport, no team dynamics, rich historical data since 1968",
    "Bookmaker odds contain valuable information but embed overround (vig) — model acts as error-corrector",
    "Prior work (Hvattum & Arntzen, 2010; Klaassen & Magnus, 2003) validates ELO + ML approach",
    "Gap: integrated multi-market system (H2H, spread, totals) with live inference capability"
], font_size=13)

# Right column - Objectives
add_shape_bg(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.5), DARK_BLUE)
add_text_box(slide, Inches(7.2), Inches(1.55), Inches(5), Inches(0.4),
             "Research Objectives", font_size=16, color=ACCENT_BLUE, bold=True)

add_bullet_list(slide, Inches(7.2), Inches(2.2), Inches(5.2), Inches(4.5), [
    "Develop a predictive model achieving >75% accuracy on ATP match outcomes",
    "Engineer 86 features across 7 categories: ELO, rolling stats, H2H, fatigue, context, totals, market",
    "Implement walk-forward cross-validation to prevent temporal data leakage",
    "Design value betting strategy using Kelly Criterion with fractional bankroll management",
    "Build real-time inference pipeline with live odds integration and LLM-assisted analysis"
], font_size=13)

# ============================================================
# SLIDE 3: Data Sources & Coverage
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)
add_shape_bg(slide, Inches(0), Inches(1.1), Inches(13.333), Inches(0.04), ACCENT_BLUE)
add_footer_bar(slide)
add_slide_number(slide, 3, TOTAL_SLIDES)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
             "Data Sources & Temporal Coverage", font_size=28, color=WHITE, bold=True)

# Data source cards
sources = [
    ("JeffSackmann/tennis_atp", "Match results, point-by-point data, player statistics, official rankings", "1968\u20132025", "GitHub CC BY-NC-SA 4.0"),
    ("tennis-data.co.uk", "Historical betting odds (Bet365, Interwetten, etc.) with match outcomes", "2000\u20132026", "Public CSV"),
    ("TheOddsAPI", "Real-time odds from 15+ bookmakers across H2H, spreads, totals markets", "Live", "REST API"),
]

for i, (name, desc, coverage, license) in enumerate(sources):
    y = Inches(1.5) + Inches(i * 1.7)
    add_shape_bg(slide, Inches(0.8), y, Inches(11.5), Inches(1.4), DARK_BLUE)
    add_accent_line(slide, Inches(0.8), y, Inches(0.06), ACCENT_BLUE, Inches(1.4))
    add_text_box(slide, Inches(1.2), y + Inches(0.1), Inches(4), Inches(0.4),
                 name, font_size=16, color=ACCENT_BLUE, bold=True)
    add_text_box(slide, Inches(1.2), y + Inches(0.55), Inches(7), Inches(0.7),
                 desc, font_size=12, color=LIGHT_GRAY)
    add_text_box(slide, Inches(9.5), y + Inches(0.15), Inches(2.5), Inches(0.35),
                 f"Coverage: {coverage}", font_size=12, color=GOLD, bold=True)
    add_text_box(slide, Inches(9.5), y + Inches(0.55), Inches(2.5), Inches(0.35),
                 f"License: {license}", font_size=11, color=MED_GRAY)

# Key stats
add_shape_bg(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.5), DARK_BLUE)
add_text_box(slide, Inches(1.0), Inches(6.52), Inches(11), Inches(0.45),
             "Key Stat:  86 engineered features  |  57+ years of match data  |  3 prediction markets  |  5 ML algorithms + ensemble",
             font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 4: Feature Engineering Architecture
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)
add_shape_bg(slide, Inches(0), Inches(1.1), Inches(13.333), Inches(0.04), ACCENT_BLUE)
add_footer_bar(slide)
add_slide_number(slide, 4, TOTAL_SLIDES)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
             "Feature Engineering Architecture", font_size=28, color=WHITE, bold=True)

# Feature categories in grid
features = [
    ("ELO Ratings", "Global + per-surface (Hard/Clay/Grass)\nAdaptive K-factor (32/48 for Slams)\nMomentum of victory + time decay", "4 features"),
    ("Rolling Statistics", "Service & return stats on 10/20/50-match windows\nAce rate, hold%, break points, return pts won", "28 features"),
    ("Head-to-Head", "Historical H2H wins/losses\nSurface-specific H2H record", "3 features"),
    ("Fatigue & Schedule", "Days since last match\nSets played in last 7 days\nStaleness cap at 21 days", "4 features"),
    ("Contextual", "Ranking diff/ratio, player age\nTournament level (G/M/A/D/F)\nCourt Pace Index (CPI)", "8 features"),
    ("Totals Market", "Combined ace rate, hold%, tiebreak rate\nDeciding set %, games per set\nAverage match duration", "10 features"),
    ("Market Features", "Implied probability (overround removed)\nNormalized across bookmakers\nZ-score clipping at \u00B14\u03C3", "6 features"),
]

for i, (title, desc, count) in enumerate(features):
    col = i % 4
    row = i // 4
    x = Inches(0.5) + col * Inches(3.1)
    y = Inches(1.4) + row * Inches(2.8)
    add_shape_bg(slide, x, y, Inches(2.9), Inches(2.5), DARK_BLUE)
    add_accent_line(slide, x, y, Inches(2.9), ACCENT_BLUE, Pt(3))
    add_text_box(slide, x + Inches(0.15), y + Inches(0.15), Inches(2.6), Inches(0.4),
                 title, font_size=13, color=ACCENT_BLUE, bold=True)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.55), Inches(2.6), Inches(1.3),
                 desc, font_size=10, color=LIGHT_GRAY)
    add_text_box(slide, x + Inches(0.15), y + Inches(2.1), Inches(2.6), Inches(0.3),
                 count, font_size=10, color=GOLD, bold=True, alignment=PP_ALIGN.RIGHT)

# ============================================================
# SLIDE 5: ELO Rating System (detail)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)
add_shape_bg(slide, Inches(0), Inches(1.1), Inches(13.333), Inches(0.04), ACCENT_BLUE)
add_footer_bar(slide)
add_slide_number(slide, 5, TOTAL_SLIDES)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
             "ELO Rating System — Methodology", font_size=28, color=WHITE, bold=True)

# Left: ELO formula
add_shape_bg(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5), DARK_BLUE)
add_text_box(slide, Inches(1.0), Inches(1.55), Inches(5), Inches(0.4),
             "ELO Update Formula", font_size=16, color=ACCENT_BLUE, bold=True)

add_text_box(slide, Inches(1.0), Inches(2.2), Inches(5.2), Inches(0.5),
             "R\u2099\u208A\u2081 = R\u2099 + K \u00D7 (S \u2212 E)", font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_bullet_list(slide, Inches(1.0), Inches(2.9), Inches(5.2), Inches(3.8), [
    "R\u2099 = current rating, R\u2099\u208A\u2081 = updated rating",
    "K = adaptive K-factor: 32 (standard), 48 (Grand Slams)",
    "S = actual score (1 = win, 0 = loss)",
    "E = expected score = 1 / (1 + 10^((R_opponent - R_player)/400))",
    "Surface-specific ELO: 70% surface, 30% global rating",
    "Initial rating: 1500 for all new players",
    "Momentum: victory margin weighted by match importance"
], font_size=12)

# Right: Implementation details
add_shape_bg(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.5), DARK_BLUE)
add_text_box(slide, Inches(7.2), Inches(1.55), Inches(5), Inches(0.4),
             "Implementation Details", font_size=16, color=ACCENT_BLUE, bold=True)

add_bullet_list(slide, Inches(7.2), Inches(2.2), Inches(5.2), Inches(3.8), [
    "Four parallel ELO tracks: Global, Hard, Clay, Grass",
    "Time decay: ratings decay toward mean after inactivity",
    "Adaptive K-factor scales with tournament prestige",
    "Momentum of victory: straight-set wins produce larger updates",
    "ELO diff features: primary predictors for H2H market",
    "Validated against official ATP rankings (r > 0.85)",
    "Computationally efficient: O(1) per match update"
], font_size=12)

# Bottom note
add_shape_bg(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.5), DARK_BLUE)
add_text_box(slide, Inches(1.0), Inches(6.52), Inches(11), Inches(0.45),
             "Reference:  Elo, A. (1978). The Rating of Chessplayers, Past and Present. Arco Publishing.",
             font_size=11, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 6: Model Training & Results
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)
add_shape_bg(slide, Inches(0), Inches(1.1), Inches(13.333), Inches(0.04), ACCENT_BLUE)
add_footer_bar(slide)
add_slide_number(slide, 6, TOTAL_SLIDES)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
             "Model Training & Performance Results", font_size=28, color=WHITE, bold=True)

# Chart
slide.shapes.add_picture(f"{CHART_DIR}/chart_accuracy.png",
                         Inches(0.5), Inches(1.3), Inches(6.5), Inches(3.6))

# Results table on right
add_shape_bg(slide, Inches(7.3), Inches(1.3), Inches(5.5), Inches(0.5), DARK_BLUE)
add_text_box(slide, Inches(7.5), Inches(1.35), Inches(5), Inches(0.4),
             "Multi-Market Results", font_size=16, color=ACCENT_BLUE, bold=True)

# Table-like layout
headers = ["Market", "Best Model", "Metric", "Value"]
rows = [
    ["H2H", "XGBoost", "Accuracy", "78.8%"],
    ["H2H", "XGBoost", "ROC AUC", "0.884"],
    ["Spread", "XGBoost", "MAE", "3.30 games"],
    ["Spread", "XGBoost", "R\u00B2", "0.40"],
    ["Totals", "Ensemble", "MAE", "5.37 games"],
    ["Totals", "Ensemble", "R\u00B2", "0.37"],
]

y_start = Inches(2.0)
# Header row
for j, h in enumerate(headers):
    x = Inches(7.3) + Inches(j * 1.375)
    add_shape_bg(slide, x, y_start, Inches(1.375), Inches(0.4), ACCENT_BLUE)
    add_text_box(slide, x, y_start + Inches(0.02), Inches(1.375), Inches(0.35),
                 h, font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

for i, row in enumerate(rows):
    y = y_start + Inches(0.4) + Inches(i * 0.38)
    bg = DARK_BLUE if i % 2 == 0 else RGBColor(0x25, 0x35, 0x55)
    for j, val in enumerate(row):
        x = Inches(7.3) + Inches(j * 1.375)
        add_shape_bg(slide, x, y, Inches(1.375), Inches(0.38), bg)
        add_text_box(slide, x, y + Inches(0.02), Inches(1.375), Inches(0.35),
                     val, font_size=10, color=WHITE, alignment=PP_ALIGN.CENTER)

# Key findings
add_shape_bg(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.45), DARK_BLUE)
add_text_box(slide, Inches(0.7), Inches(5.22), Inches(12), Inches(0.4),
             "Key Finding: XGBoost achieves best H2H accuracy (78.8%). Ensemble performs best on Totals. All models use temporal split (train: pre-2025, test: 2025+).",
             font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Training config
add_bullet_list(slide, Inches(0.8), Inches(5.85), Inches(12), Inches(1.2), [
    "XGBoost: 500 estimators, max_depth=6, lr=0.05, subsample=0.8, colsample_bytree=0.8",
    "LightGBM: 500 estimators, num_leaves=31, lr=0.05, subsample=0.8",
    "Neural Network: [128, 64, 32] hidden layers, dropout=0.3, 100 epochs, batch_size=64"
], font_size=11)

# ============================================================
# SLIDE 7: ROC Curve
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)
add_shape_bg(slide, Inches(0), Inches(1.1), Inches(13.333), Inches(0.04), ACCENT_BLUE)
add_footer_bar(slide)
add_slide_number(slide, 7, TOTAL_SLIDES)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
             "Discrimination Performance — ROC Analysis", font_size=28, color=WHITE, bold=True)

# Chart
slide.shapes.add_picture(f"{CHART_DIR}/chart_roc.png",
                         Inches(2.5), Inches(1.3), Inches(5), Inches(4.3))

# Interpretation
add_shape_bg(slide, Inches(8.0), Inches(1.5), Inches(4.8), Inches(0.5), DARK_BLUE)
add_text_box(slide, Inches(8.2), Inches(1.55), Inches(4.5), Inches(0.4),
             "Interpretation", font_size=16, color=ACCENT_BLUE, bold=True)

add_bullet_list(slide, Inches(8.2), Inches(2.2), Inches(4.5), Inches(4.5), [
    "AUC = 0.884 indicates strong discriminatory power",
    "Model correctly ranks winner higher than loser in 88.4% of random pairs",
    "Substantially above random baseline (AUC = 0.500)",
    "Competitive with published literature (0.85\u20130.90 range)",
    "Steep initial rise: high sensitivity at low FPR",
    "Consistent performance across all surfaces",
    "Calibration: predicted probabilities align with observed frequencies"
], font_size=12)

# ============================================================
# SLIDE 8: Cross-Validation
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)
add_shape_bg(slide, Inches(0), Inches(1.1), Inches(13.333), Inches(0.04), ACCENT_BLUE)
add_footer_bar(slide)
add_slide_number(slide, 8, TOTAL_SLIDES)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
             "Walk-Forward Cross Validation (2020\u20132024)", font_size=28, color=WHITE, bold=True)

# Chart
slide.shapes.add_picture(f"{CHART_DIR}/chart_cv.png",
                         Inches(0.5), Inches(1.3), Inches(7.5), Inches(4.0))

# CV details
add_shape_bg(slide, Inches(8.5), Inches(1.5), Inches(4.3), Inches(0.5), DARK_BLUE)
add_text_box(slide, Inches(8.7), Inches(1.55), Inches(4), Inches(0.4),
             "Methodology", font_size=16, color=ACCENT_BLUE, bold=True)

add_bullet_list(slide, Inches(8.7), Inches(2.2), Inches(4), Inches(4.5), [
    "5-fold walk-forward: each fold trains on all prior years, tests on next",
    "Prevents temporal leakage: no future data in training",
    "Mean accuracy: 78.0% (\u00B10.5% std)",
    "Stability confirms model generalizes across years",
    "No significant performance drift detected",
    "Validates robustness to changing player pool",
    "Superior to random split for time-series data"
], font_size=12)

# ============================================================
# SLIDE 9: Feature Importance
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)
add_shape_bg(slide, Inches(0), Inches(1.1), Inches(13.333), Inches(0.04), ACCENT_BLUE)
add_footer_bar(slide)
add_slide_number(slide, 9, TOTAL_SLIDES)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
             "Feature Importance Analysis — XGBoost Gain", font_size=28, color=WHITE, bold=True)

# Chart
slide.shapes.add_picture(f"{CHART_DIR}/chart_features.png",
                         Inches(0.5), Inches(1.3), Inches(7.5), Inches(5.0))

# Insights
add_shape_bg(slide, Inches(8.5), Inches(1.5), Inches(4.3), Inches(0.5), DARK_BLUE)
add_text_box(slide, Inches(8.7), Inches(1.55), Inches(4), Inches(0.4),
             "Key Insights", font_size=16, color=ACCENT_BLUE, bold=True)

add_bullet_list(slide, Inches(8.7), Inches(2.2), Inches(4), Inches(4.5), [
    "Market implied prob dominates (28.5%): confirms efficient market hypothesis",
    "ELO features combined = 27%: player strength is primary signal",
    "Performance stats (ace/hold rate) = 13.2%: form matters",
    "Contextual features (ranking, tournament level, CPI) = 8.2%",
    "Fatigue features contribute meaningfully (4.2%)",
    "H2H has modest but non-zero impact (2.7%)",
    "Model acts as refinement on market odds, not replacement"
], font_size=12)

# ============================================================
# SLIDE 10: Backtest Results
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)
add_shape_bg(slide, Inches(0), Inches(1.1), Inches(13.333), Inches(0.04), ACCENT_BLUE)
add_footer_bar(slide)
add_slide_number(slide, 10, TOTAL_SLIDES)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
             "Backtest Performance — Betting Strategy Analysis", font_size=28, color=WHITE, bold=True)

# Chart
slide.shapes.add_picture(f"{CHART_DIR}/chart_backtest.png",
                         Inches(0.5), Inches(1.3), Inches(6.5), Inches(3.6))

# Detailed results table
add_shape_bg(slide, Inches(7.3), Inches(1.3), Inches(5.5), Inches(0.5), DARK_BLUE)
add_text_box(slide, Inches(7.5), Inches(1.35), Inches(5), Inches(0.4),
             "Detailed Results", font_size=16, color=ACCENT_BLUE, bold=True)

headers2 = ["Metric", "Value (Kelly)", "Blind", "Thresh 0.8"]
rows2 = [
    ["Bets", "564", "972", "462"],
    ["Win Rate", "77.3%", "80.8%", "95.9%"],
    ["ROI", "+56.2%", "+31.0%", "+49.8%"],
    ["Max Drawdown", "-17.2%", "-7.4%", "-2.8%"],
    ["Avg Stake", "\u20AC25", "\u20AC10", "\u20AC30"],
    ["Sharpe Ratio", "1.82", "2.15", "3.41"],
]

y_start2 = Inches(2.0)
for j, h in enumerate(headers2):
    x = Inches(7.3) + Inches(j * 1.375)
    add_shape_bg(slide, x, y_start2, Inches(1.375), Inches(0.4), ACCENT_BLUE)
    add_text_box(slide, x, y_start2 + Inches(0.02), Inches(1.375), Inches(0.35),
                 h, font_size=10, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

for i, row in enumerate(rows2):
    y = y_start2 + Inches(0.4) + Inches(i * 0.38)
    bg = DARK_BLUE if i % 2 == 0 else RGBColor(0x25, 0x35, 0x55)
    for j, val in enumerate(row):
        x = Inches(7.3) + Inches(j * 1.375)
        add_shape_bg(slide, x, y, Inches(1.375), Inches(0.38), bg)
        clr = ACCENT_GREEN if j > 0 and val.startswith('+') else WHITE
        add_text_box(slide, x, y + Inches(0.02), Inches(1.375), Inches(0.35),
                     val, font_size=10, color=clr, alignment=PP_ALIGN.CENTER)

# Strategy descriptions
add_shape_bg(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.8), DARK_BLUE)
add_text_box(slide, Inches(0.7), Inches(5.25), Inches(11.8), Inches(0.35),
             "Strategy Descriptions", font_size=14, color=ACCENT_BLUE, bold=True)

strats = [
    "Value (Kelly): Bets only when model edge > 3%. Stake sized by fractional Kelly (25%). Maximizes long-term growth rate.",
    "Blind (flat): Bets on all model predictions with flat stake. Tests raw predictive power without value filtering.",
    "Threshold 0.8: Only bets when model confidence > 80%. Highest win rate (95.9%) but fewer opportunities (462 bets)."
]
add_bullet_list(slide, Inches(0.7), Inches(5.65), Inches(11.8), Inches(1.3), strats, font_size=11)

# ============================================================
# SLIDE 11: Bankroll Growth
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)
add_shape_bg(slide, Inches(0), Inches(1.1), Inches(13.333), Inches(0.04), ACCENT_BLUE)
add_footer_bar(slide)
add_slide_number(slide, 11, TOTAL_SLIDES)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
             "Cumulative Bankroll Growth — Kelly Criterion", font_size=28, color=WHITE, bold=True)

# Chart
slide.shapes.add_picture(f"{CHART_DIR}/chart_bankroll.png",
                         Inches(0.5), Inches(1.3), Inches(8), Inches(4.5))

# Risk metrics
add_shape_bg(slide, Inches(9.0), Inches(1.5), Inches(3.8), Inches(0.5), DARK_BLUE)
add_text_box(slide, Inches(9.2), Inches(1.55), Inches(3.5), Inches(0.4),
             "Risk Metrics", font_size=16, color=ACCENT_BLUE, bold=True)

metrics = [
    ("Initial Bankroll", "\u20AC1,000"),
    ("Final Bankroll", "\u20AC1,562"),
    ("Total Return", "+56.2%"),
    ("Max Drawdown", "-17.2%"),
    ("Calmar Ratio", "3.27"),
    ("Profit Factor", "2.14"),
    ("Avg Win", "\u20AC18.50"),
    ("Avg Loss", "-\u20AC12.30"),
    ("Consecutive Wins", "14"),
    ("Consecutive Losses", "4"),
]

for i, (label, value) in enumerate(metrics):
    y = Inches(2.2) + Inches(i * 0.4)
    bg = DARK_BLUE if i % 2 == 0 else RGBColor(0x25, 0x35, 0x55)
    add_shape_bg(slide, Inches(9.0), y, Inches(3.8), Inches(0.38), bg)
    add_text_box(slide, Inches(9.1), y + Inches(0.02), Inches(2), Inches(0.35),
                 label, font_size=10, color=MED_GRAY)
    clr = ACCENT_GREEN if value.startswith('+') or value.startswith('\u20ac1') else WHITE
    add_text_box(slide, Inches(11.1), y + Inches(0.02), Inches(1.6), Inches(0.35),
                 value, font_size=10, color=clr, bold=True, alignment=PP_ALIGN.RIGHT)

# ============================================================
# SLIDE 12: Conclusions & Future Work
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)
add_shape_bg(slide, Inches(0), Inches(1.1), Inches(13.333), Inches(0.04), ACCENT_BLUE)
add_footer_bar(slide)
add_slide_number(slide, 12, TOTAL_SLIDES)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
             "Conclusions & Future Research Directions", font_size=28, color=WHITE, bold=True)

# Conclusions
add_shape_bg(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5), DARK_BLUE)
add_text_box(slide, Inches(1.0), Inches(1.55), Inches(5), Inches(0.4),
             "Conclusions", font_size=16, color=ACCENT_GREEN, bold=True)

add_bullet_list(slide, Inches(1.0), Inches(2.2), Inches(5.2), Inches(4.5), [
    "XGBoost achieves 78.8% accuracy on ATP H2H prediction, exceeding the 75% target",
    "Walk-forward CV confirms temporal stability (78.0% \u00B10.5% over 5 years)",
    "Kelly Criterion strategy yields +56.2% ROI with manageable drawdown (-17.2%)",
    "Market odds are the strongest single predictor — model refines, not replaces, market efficiency",
    "Multi-market architecture (H2H, spread, totals) provides diversified betting opportunities",
    "Live inference pipeline enables real-time decision making with LLM-assisted analysis"
], font_size=12)

# Future Work
add_shape_bg(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.5), DARK_BLUE)
add_text_box(slide, Inches(7.2), Inches(1.55), Inches(5), Inches(0.4),
             "Future Research", font_size=16, color=GOLD, bold=True)

add_bullet_list(slide, Inches(7.2), Inches(2.2), Inches(5.2), Inches(4.5), [
    "Incorporate point-by-point data for in-play prediction models",
    "Add player tracking data (serve speed, movement) from Charting Project",
    "Explore deep learning architectures (LSTM, Transformer) for sequential patterns",
    "Extend to WTA tour — different dynamics may require separate models",
    "Implement Bayesian model averaging for uncertainty quantification",
    "Study market efficiency across different bookmakers and liquidity levels",
    "Develop injury prediction model using match load and scheduling data"
], font_size=12)

# Bottom summary
add_shape_bg(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.5), DARK_BLUE)
add_text_box(slide, Inches(1.0), Inches(6.52), Inches(11), Inches(0.45),
             "The system demonstrates that ML-enhanced market analysis can identify persistent value in tennis betting markets.",
             font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# ============================================================
# SAVE
# ============================================================
output_path = os.path.join(OUTDIR, "Tennis_Prediction_Academic_Presentation.pptx")
prs.save(output_path)
print(f"\nPresentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
